import os
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches

from app.config import get_settings
from app.db import insert_report


def build_report(source_id: int, source_name: str, insight: str, chart_paths: list[str]) -> tuple[str, str]:
    s = get_settings()
    report_dir = os.path.join(s.output_dir, "reports")
    ppt_dir = os.path.join(s.output_dir, "ppt")
    os.makedirs(report_dir, exist_ok=True)
    os.makedirs(ppt_dir, exist_ok=True)

    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    md_path = os.path.join(report_dir, f"source_{source_id}_{ts}.md")
    ppt_path = os.path.join(ppt_dir, f"source_{source_id}_{ts}.pptx")

    md = [
        f"# 网站变更监控报告 - {source_name}",
        "",
        f"- 生成时间：{datetime.now().isoformat(sep=' ', timespec='seconds')}",
        f"- 数据源ID：{source_id}",
        "",
        "## 分析洞察",
        insight,
        "",
        "## 图表文件",
    ]
    md.extend([f"- {p}" for p in chart_paths])
    with open(md_path, "w", encoding="utf-8") as f:
        f.write("\n".join(md))

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "网站变更监控自动报告"
    slide.placeholders[1].text = f"{source_name}\n{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "分析结论"
    slide2.placeholders[1].text = insight[:1200]

    for p in chart_paths[:3]:
        sld = prs.slides.add_slide(prs.slide_layouts[5])
        sld.shapes.title.text = os.path.basename(p)
        if os.path.exists(p):
            sld.shapes.add_picture(p, Inches(0.8), Inches(1.3), width=Inches(8.5))

    prs.save(ppt_path)
    insert_report(source_id, md_path, ppt_path)
    return md_path, ppt_path
